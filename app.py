from flask import Flask, request, send_file, jsonify, render_template
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os

app = Flask(__name__)
CORS(app)

# Configuración de rutas
app.template_folder = os.path.join(os.path.dirname(__file__), "templates")
app.static_folder = os.path.join(os.path.dirname(__file__), "static")

PPTX_TEMPLATE = os.path.join(os.path.dirname(__file__), "Certificado Diploma Premio.pptx")
IMAGEN_ESTRELLA = os.path.join(os.path.dirname(__file__), "static", "estrella.png")

# Ruta principal
@app.route("/")
def index():
    return render_template("index.html")

@app.route("/generar-diploma", methods=["POST"])
def generar_diploma():
    try:
        data = request.json
        print("Recibido JSON:", data)

        nombre = data.get("nombre", "Nombre")
        estrellas = data.get("estrellas", 0)
        mes = data.get("mes", "Mes")
        fecha = data.get("fecha", "Fecha")

        if estrellas is None or estrellas == "":
            estrellas = 0  # Asegurar que siempre sea un número válido
        estrellas = int(estrellas)

        print(f"Nombre: {nombre}, Estrellas: {estrellas}, Mes: {mes}, Fecha: {fecha}")

        # Verificar la existencia de la plantilla y la imagen de la estrella
        if not os.path.exists(PPTX_TEMPLATE):
            print("❌ ERROR: El archivo PPTX no se encontró")
            return jsonify({"error": "El archivo PPTX no se encontró"}), 500

        if not os.path.exists(IMAGEN_ESTRELLA):
            print("❌ ERROR: La imagen de la estrella no se encontró")
            return jsonify({"error": "La imagen de la estrella no se encontró"}), 500

        prs = Presentation(PPTX_TEMPLATE)

        # Reemplazo de texto en la presentación
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if "[Nombre]" in shape.text:
                        shape.text = shape.text.replace("[Nombre]", nombre)

                        # Aplicar estilos al nombre
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if nombre in run.text:
                                    run.font.name = "Verdana"
                                    run.font.size = Pt(36)
                            paragraph.alignment = PP_ALIGN.CENTER

                    if "[Mes]" in shape.text:
                        shape.text = shape.text.replace("[Mes]", mes)

                        # Aplicar estilos al mes
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER
                            for run in paragraph.runs:
                                if mes in run.text:
                                    run.font.name = "Tahoma"
                                    run.font.size = Pt(16)
                                    

                    if "[Fecha]" in shape.text:
                        shape.text = shape.text.replace("[Fecha]", fecha)
                        # Aplicar estilos a la fecha
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER
                            for run in paragraph.runs:
                                if fecha in run.text:
                                    run.font.name = "Tahoma"
                                    run.font.size = Pt(16)

                    if "[Cartera]" in shape.text:
                        shape.text = shape.text.replace("[Cartera]", "Banco de Occidente")
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER
                            for run in paragraph.runs:
                                if fecha in run.text:
                                    run.font.name = "Tahoma"
                                    run.font.size = Pt(16)
                                    run.font.bold = True

        # Reemplazo de estrellas con imágenes
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and "[ESTRELLAS]" in shape.text:
                    shape.text = ""  # Eliminar el marcador de estrellas

                    # Posición y tamaño de las estrellas
                    left, top, width = shape.left, shape.top, shape.width
                    estrella_width, estrella_height = Pt(18), Pt(18)  # Tamaño en puntos
                    espacio_entre_estrellas = Pt(10)
                    espacio_total = (estrella_width * estrellas) + (espacio_entre_estrellas * (estrellas - 1))
                    left_inicio = left + (width - espacio_total) / 2

                    # Insertar las imágenes de las estrellas
                    for i in range(estrellas):
                        estrella_left = left_inicio + (i * (estrella_width + espacio_entre_estrellas))
                        slide.shapes.add_picture(
                            IMAGEN_ESTRELLA,
                            estrella_left,
                            top,
                            width=estrella_width,
                            height=estrella_height
                        )

        output_filename = f"diploma_{nombre}.pptx"
        prs.save(output_filename)

        print("✅ Diploma generado correctamente:", output_filename)
        return send_file(output_filename, as_attachment=True)

    except Exception as e:
        print("❌ ERROR en Flask:", str(e))
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(debug=True)
